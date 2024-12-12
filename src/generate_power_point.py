from datetime import datetime

import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.util import Inches
from pptx.util import Pt


def create_table_df(df):
    table = []

    for ticker in df["TICKER"].unique():
        stock_data = df[df["TICKER"] == ticker]
        start_price = stock_data.iloc[0]["OPEN"]
        end_price = stock_data.iloc[-1]["CLOSE"]
        n_days = len(df['DATETIME'])
        returns = ((end_price - start_price) / start_price) * 100
        volatility = stock_data["CLOSE"].pct_change().std() * np.sqrt(n_days) * 100

        table.append({
            "Ticker": ticker,
            "Start Price": f"{start_price:.2f}",
            "End Price": f"{end_price:.2f}",
            "Return (%)": f"{returns:.2f}",
            "Volatility (%)": f"{volatility:.2f}"
        })

    table_df = pd.DataFrame(table)

    return table_df


def generate_ppt(df, output_file="stock_report.pptx"):
    table_df = create_table_df(df)

    ppt = Presentation("slides/template.pptx")

    # Title slide
    title_slide = ppt.slides[0]
    title_slide.shapes.title.text = "Python generated report (Demo)"
    idxs = [shape.placeholder_format.idx for shape in title_slide.placeholders]
    title_slide.placeholders[idxs[-1]].text = "Stock Report"

    # Table slide
    create_table_slide(ppt, table_df)

    # Bar chart slide
    create_bar_chart_slide(ppt, df)

    # Line chart slide
    create_line_chart_slide(ppt, df)

    # End slide
    end_slide = ppt.slides.add_slide(ppt.slide_layouts[0])
    end_slide.shapes.title.text = "Made by Henrik A. Kjeldsberg"

    ppt.save(output_file)
    return output_file


def create_bar_chart_slide(ppt, df):
    bar_chart_slide_layout = ppt.slide_layouts[9]
    slide = ppt.slides.add_slide(bar_chart_slide_layout)
    title = slide.shapes.title
    title.text = "Return per stock statistics"

    # Add chart data
    chart_data = CategoryChartData()
    chart_data.categories = df['TICKER'].unique()
    returns_by_ticker = df.groupby("TICKER")["RETURN"].mean()
    chart_data.add_series('Returns', returns_by_ticker.values)

    # Add chart to slide
    x, y, cx, cy = Inches(1), Inches(1.5), Inches(10), Inches(5)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
    ).chart

    # Style the chart
    chart.has_legend = False


def create_line_chart_slide(ppt, df):
    line_chart_slide_layout = ppt.slide_layouts[9]
    slide = ppt.slides.add_slide(line_chart_slide_layout)
    title = slide.shapes.title
    title.text = "Opening price per stock"

    # Add chart data
    chart_data = CategoryChartData()
    chart_data.categories = df['DATETIME'].dt.tz_localize(None).unique()
    for ticker in df['TICKER'].unique():
        values = df[df['TICKER'] == ticker]
        line = chart_data.add_series(ticker, values['OPEN'].to_list())

    # Add chart
    x, y, cx, cy = Inches(.5), Inches(1.5), Inches(12), Inches(5.5)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data
    ).chart

    # Manualmcoloring
    colors = [
        RGBColor(68, 1, 84),
        RGBColor(72, 35, 116),
        RGBColor(64, 67, 135),
        RGBColor(52, 94, 141),
        RGBColor(41, 120, 142),
        RGBColor(32, 144, 140),
        RGBColor(53, 173, 129),
        RGBColor(109, 199, 98),
        RGBColor(180, 222, 44),
        RGBColor(253, 231, 37),
    ]
    for i, series in enumerate(chart.series):
        series.format.line.color.rgb = colors[i]

    # Styling
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.TOP
    chart.legend.include_in_layout = False


def create_table_slide(ppt, df):
    table_slide_layout = ppt.slide_layouts[9]
    slide = ppt.slides.add_slide(table_slide_layout)
    title = slide.shapes.title
    title.text = "Return per stock statistics"

    # Create a table with the statistics
    rows, cols = df.shape
    left, top, right, bottom = Inches(2), Inches(2), Inches(9), Inches(4)
    table = slide.shapes.add_table(rows + 1, cols, left, top, right, bottom).table

    def style_table(table, df):
        for col_idx, col_name in enumerate(df.columns):
            cell = table.cell(0, col_idx)
            cell.text = col_name
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0, 102, 204)  # Blue background

            cell.text_frame.paragraphs[0].font.bold = True
            cell.text_frame.paragraphs[0].font.size = Pt(12)
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # White text
            cell.text_frame.paragraphs[0].alignment = 1  # Center alignment

        for row_idx, row in enumerate(df.values):
            for col_idx, value in enumerate(row):
                cell = table.cell(row_idx + 1, col_idx)
                cell.text = str(value)
                cell.text_frame.paragraphs[0].font.size = Pt(10)
                cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)  # Black text
                cell.text_frame.paragraphs[0].alignment = 1  # Center alignment

                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(240, 240, 240)  # Light gray background

    style_table(table, df)

    # Set the table style
    table.first_row = True  # Mark first row as a header row
    table.horz_banding = False  # Turn off horizontal banding
    table.vert_banding = False  # Turn off vertical banding

    for col_idx, column_name in enumerate(df.columns):
        table.cell(0, col_idx).text = column_name.capitalize()

    for row_idx, row_data in enumerate(df.values, start=1):
        for col_idx, value in enumerate(row_data):
            if isinstance(value, float):
                value = f"{float(value):.2f}"
            elif isinstance(value, datetime):
                value = value.strftime("%m/%d/%Y")
            else:
                value = str(value)
            table.cell(row_idx, col_idx).text = value


def main():
    slide = ppt = df = feature = output_file = None
    slide.shapes.title.text = "Test zero"
    slide.shapes.textbox.text = "Test zero"
    slide.placeholders[0].text = "Generated XXX"
    # print(slide)
    # slide.placeholders[0].text ="Test 2"
    # slide.placeholders[1].text ="Test 2"
    #
    # title_slide_layout = ppt.slide_layouts[0]
    # slide = ppt.slides.add_slide(title_slide_layout)
    # slide.shapes.title.text = "Dynamic Report"
    # #slide.placeholders[1].text = "Generated using  Streamlit + Python + python-pptx"

    # Page 2 - Raw data
    table_slide_layout = ppt.slide_layouts[5]
    slide = ppt.slides.add_slide(table_slide_layout)
    title = slide.shapes.title
    title.text = "Data Overview"

    rows, cols = df.shape
    max_rows = 1
    table = slide.shapes.add_table(max_rows + 1, cols, Inches(0.2), Inches(0.2), Inches(9), Inches(4)).table

    for col_idx, column_name in enumerate(df.columns):
        table.cell(0, col_idx).text = column_name

    for row_idx, row_data in enumerate(df.head(max_rows).values, start=1):
        for col_idx, value in enumerate(row_data):
            if isinstance(value, float):
                value = f"{float(value):.2f}"
            elif isinstance(value, datetime):
                value = value.strftime("%m/%d/%Y")
            else:
                value = str(value)
            table.cell(row_idx, col_idx).text = value

    # Page 3: Chart
    if not df.empty:
        fig, ax = plt.subplots()
        df.groupby("TICKER")[feature].plot(ax=ax)

        plt.title(f"{feature.capitalize()} time series")
        plt.legend(loc="upper left")

        chart_image = "slides/chart.png"
        plt.savefig(chart_image)
        plt.close()

        chart_slide_layout = ppt.slide_layouts[5]
        slide = ppt.slides.add_slide(chart_slide_layout)
        title = slide.shapes.title
        title.text = "Chart Representation"
        slide.shapes.add_picture(chart_image, Inches(1), Inches(1.5), Inches(7), Inches(5))

    ppt.save(output_file)
    return output_file
